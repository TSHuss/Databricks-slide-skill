#!/usr/bin/env python3
"""
Databricks Slide Deck Generator (Template-Based)

Generates PowerPoint presentations using the official Databricks corporate template.
Output can be imported directly into Google Slides.

Usage:
    python generate-pptx.py --input content.json --output presentation.pptx
"""

import json
import argparse
import re
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"Error: python-pptx is required. Install with: pip3 install python-pptx")
    print(f"Details: {e}")
    sys.exit(1)

# =============================================================================
# Constants
# =============================================================================

SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
TEMPLATE_PATH = SKILL_DIR / "assets" / "databricks" / "template.pptx"
THEME_PATH = SKILL_DIR / "themes" / "databricks.json"

# Valid slide types (25 total)
VALID_SLIDE_TYPES = {
    # Existing types (17)
    "title", "section", "content", "two-column", "three-column",
    "big-number", "callout", "quote", "closing",
    "agenda", "timeline", "icon-grid", "stat-row", "pros-cons",
    "comparison", "checklist", "logos",
    # New types (8)
    "two-column-icons", "three-column-icons", "cards",
    "card-right", "card-left", "card-full",
    "one-column", "section-description"
}

# Layout name mappings (our type -> template layout name patterns)
# For slides with prefer_dark=True, get_layout searches dark_layouts first
LAYOUT_MAPPINGS = {
    # Structural slides (dark by default) - use Databricks dark templates
    "title": ["1_3 Title Slide B - Dark", "3 Title Slide B - Light", "TITLE"],
    "section": ["Content E - Power Statement 3", "SECTION_HEADER"],  # Dark statement layout
    "callout": ["Content E - Power Statement 2_1", "MAIN_POINT"],  # Dark with title+subtitle
    "quote": ["Content E - Power Statement 2_1", "MAIN_POINT"],  # Dark with title+subtitle
    "closing": ["Z - Closing Dark", "Z - Closing Light"],
    # Content slides (light by default)
    "content": ["7 Content A - Basic", "TITLE_AND_BODY"],
    "two-column": ["9 Content B - 2 Column", "TITLE_AND_TWO_COLUMNS"],
    "three-column": ["11 Content C - 3 Column"],
    "big-number": ["Content E - Power Statement 1", "BIG_NUMBER"],
    # New template types (these are already on Master 1/2)
    "two-column-icons": ["10 Content B - 2 Column w/ Icon Spot"],
    "three-column-icons": ["12 Content C - 3 Column w/ Icon Spot"],
    "cards": ["13 Content C - 3 Column Cards"],
    "card-right": ["14 Content D - Card Right"],
    "card-left": ["15 Content D - Card Left"],
    "card-full": ["16 Content D - Card Large"],
    "one-column": ["7 Content A - Basic", "ONE_COLUMN_TEXT"],  # Master 1/2 for footer
    "section-description": ["Content E - Power Statement 2", "SECTION_TITLE_AND_DESCRIPTION"],  # Master 1 for footer
    # Hybrid types (use CUSTOM for clean slate with footer from master)
    "agenda": ["CUSTOM"],
    "timeline": ["CUSTOM"],
    "icon-grid": ["CUSTOM"],
    "stat-row": ["CUSTOM"],
    "pros-cons": ["CUSTOM"],
    "comparison": ["CUSTOM"],
    "checklist": ["CUSTOM"],
    "logos": ["CUSTOM"],
}


def load_colors_from_theme(theme_path: Path = THEME_PATH) -> Dict[str, str]:
    """Load color palette from theme JSON file.

    Falls back to hardcoded Databricks brand colors if theme file
    is missing or invalid.
    """
    try:
        with open(theme_path, 'r') as f:
            theme = json.load(f)

        return {
            "accent": theme["modes"]["light"]["accent"],
            "dark_bg": theme["modes"]["dark"]["background"],
            "light_bg": theme["modes"]["light"]["background"],
            "text_dark": theme["modes"]["light"]["text_primary"],
            "text_light": theme["modes"]["dark"]["text_primary"],
            "text_secondary": theme["modes"]["light"]["text_secondary"],
            "green": theme["elements"]["pros_header_color"],
            "red": theme["elements"]["cons_header_color"],
            "divider": theme["elements"]["stat_row_divider"],
        }
    except FileNotFoundError:
        print(f"Warning: Theme file not found at {theme_path}, using defaults")
    except json.JSONDecodeError as e:
        print(f"Warning: Invalid JSON in theme file: {e}, using defaults")
    except KeyError as e:
        print(f"Warning: Missing key in theme file: {e}, using defaults")

    # Fallback to hardcoded Databricks brand colors
    return {
        "accent": "#FF3621",
        "dark_bg": "#1B3139",
        "light_bg": "#F5F3F0",
        "text_dark": "#1B3139",
        "text_light": "#FFFFFF",
        "text_secondary": "#6B7280",
        "green": "#10B981",
        "red": "#EF4444",
        "divider": "#E5E7EB",
    }


# Colors loaded from theme file (with hardcoded fallback)
COLORS = load_colors_from_theme()


def load_font_from_theme(theme_path: Path = THEME_PATH) -> str:
    """Load font family from theme JSON file."""
    try:
        with open(theme_path, 'r') as f:
            theme = json.load(f)
        return theme.get("typography", {}).get("font_family", "DM Sans")
    except (FileNotFoundError, json.JSONDecodeError, KeyError):
        return "DM Sans"


# Font family from theme (with fallback)
FONT_FAMILY = load_font_from_theme()

_color_cache: Dict[str, RGBColor] = {}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor with caching."""
    if hex_color in _color_cache:
        return _color_cache[hex_color]
    clean_hex = hex_color.lstrip('#')
    color = RGBColor(int(clean_hex[0:2], 16), int(clean_hex[2:4], 16), int(clean_hex[4:6], 16))
    _color_cache[hex_color] = color
    return color


# Precompiled regex for accent text parsing
_ACCENT_PATTERN = re.compile(r'\*([^*]+)\*')


def parse_accent_text(text: str) -> List[Tuple[str, bool]]:
    """Parse text for *accent* markers.

    Returns list of (text, is_accent) tuples.
    Example: "Hello *world* today" -> [("Hello ", False), ("world", True), (" today", False)]
    """
    if '*' not in text:
        return [(text, False)]

    segments = []
    last_end = 0

    for match in _ACCENT_PATTERN.finditer(text):
        # Add text before this match (if any)
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], False))
        # Add the accented text (without asterisks)
        segments.append((match.group(1), True))
        last_end = match.end()

    # Add remaining text after last match
    if last_end < len(text):
        segments.append((text[last_end:], False))

    return segments if segments else [(text, False)]


# =============================================================================
# Generator Class
# =============================================================================

class DatabricksSlideGenerator:
    """Generate Databricks-branded PowerPoint presentations using corporate template."""

    def __init__(self, template_path: Path = TEMPLATE_PATH):
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load template
        self.prs = Presentation(str(template_path))
        self.slide_count = 0

        # Layout caches for background mode support:
        #
        # - layouts: default cache, prefers light (first occurrence wins during iteration)
        # - light_layouts: explicit light-background layouts only (from light masters)
        # - dark_layouts: explicit dark-background layouts only (from dark masters)
        #
        # Current usage:
        #   - Template slides use `layouts` (defaults to light)
        #   - Hybrid slides use `dark_layouts` via prefer_dark=True
        #
        # Future: light_layouts enables explicit "prefer_light=True" for user-selected
        # light/dark mode per slide. Keeping all three caches for this extensibility.
        self.layouts: Dict[str, Any] = {}
        self.light_layouts: Dict[str, Any] = {}
        self.dark_layouts: Dict[str, Any] = {}

        for master in self.prs.slide_masters:
            is_dark_master = self._is_dark_background(master)

            for layout in master.slide_layouts:
                # Store in light/dark specific cache
                if is_dark_master:
                    self.dark_layouts[layout.name] = layout
                else:
                    self.light_layouts[layout.name] = layout

                # Store in main cache only if not already present (prefer first/light)
                if layout.name not in self.layouts:
                    self.layouts[layout.name] = layout

        # Clear existing slides (keep layouts)
        self._clear_slides()

    def _clear_slides(self) -> None:
        """Remove all slides from the presentation while keeping layouts."""
        # Delete slides from end to start to avoid index issues
        for i in range(len(self.prs.slides) - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def _is_dark_background(self, master) -> bool:
        """Check if a slide master has a dark background color."""
        try:
            fill = master.background.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                    rgb = fill.fore_color.rgb
                    if rgb:
                        rgb_str = str(rgb).upper()
                        # Known dark color (Databricks dark bg)
                        if rgb_str == '1B3139':
                            return True
                        # General darkness check: RGB sum < 384 means avg < 128
                        try:
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            if (r + g + b) < 384:
                                return True
                        except (ValueError, IndexError):
                            pass
        except (AttributeError, TypeError):
            pass
        return False

    def get_layout(self, slide_type: str, prefer_dark: bool = False):
        """Get the best matching layout for a slide type.

        Args:
            slide_type: The type of slide to get a layout for
            prefer_dark: If True, prefer dark background layouts when available
        """
        patterns = LAYOUT_MAPPINGS.get(slide_type, ["BLANK"])

        # Choose which caches to search based on preference
        if prefer_dark:
            search_caches = [self.dark_layouts, self.layouts]
        else:
            search_caches = [self.layouts]

        for pattern in patterns:
            for cache in search_caches:
                # Try exact match first
                if pattern in cache:
                    return cache[pattern]
            for cache in search_caches:
                # Try partial match
                for name, layout in cache.items():
                    if pattern in name:
                        return layout

        # Fallback to BLANK
        return self.layouts.get("BLANK", list(self.layouts.values())[0])

    def get_placeholder(self, slide, idx: int = None, ph_type: int = None):
        """Get placeholder by index or type."""
        for shape in slide.placeholders:
            if idx is not None and shape.placeholder_format.idx == idx:
                return shape
            if ph_type is not None and shape.placeholder_format.type == ph_type:
                return shape
        return None

    def get_placeholders_by_type(self, slide, ph_type: int) -> List:
        """Get all placeholders of a certain type, sorted by position (left to right, top to bottom)."""
        matching = []
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                matching.append(shape)
        # Sort by top position first, then left
        matching.sort(key=lambda s: (s.top, s.left))
        return matching

    def fill_text(self, placeholder, text: str, font_size: int = None,
                  bold: bool = None, color: str = None) -> None:
        """Fill a placeholder with styled text.

        Supports accent text: wrap words in *asterisks* to highlight
        them in the accent color (Databricks orange).
        """
        if placeholder is None:
            return

        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]

        # Parse for accent markers
        segments = parse_accent_text(text)
        has_accent = any(is_accent for _, is_accent in segments)

        if has_accent:
            # Use runs for mixed formatting
            for i, (segment_text, is_accent) in enumerate(segments):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                run.text = segment_text

                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold

                # Apply accent color or specified color
                if is_accent:
                    run.font.color.rgb = hex_to_rgb(COLORS["accent"])
                elif color:
                    run.font.color.rgb = hex_to_rgb(color)
        else:
            # Simple case: no accent markers
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = hex_to_rgb(color)

    def fill_bullets(self, placeholder, items: List[str], font_size: int = None) -> None:
        """Fill a placeholder with bullet points."""
        if placeholder is None or not items:
            return

        tf = placeholder.text_frame
        tf.clear()

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            if font_size:
                p.font.size = Pt(font_size)

    def add_textbox(self, slide, text: str, left: float, top: float,
                    width: float, height: float, font_size: int = 18,
                    bold: bool = False, color: str = None,
                    alignment: int = None) -> None:
        """Add a styled textbox to a slide."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = hex_to_rgb(color)
        if alignment:
            p.alignment = alignment

    def _create_slide(self, slide_type: str, data: Dict[str, Any], prefer_dark: bool = False):
        """Create slide with layout and handle common setup.

        Handles:
        - Incrementing slide count
        - Getting the appropriate layout (light or dark)
        - Adding the slide to the presentation
        - Adding speaker notes if present in data

        Args:
            slide_type: The type of slide (maps to LAYOUT_MAPPINGS)
            data: Slide data dict, may contain "notes" key
            prefer_dark: If True, prefer dark-background layout

        Returns:
            The created slide object
        """
        self.slide_count += 1
        layout = self.get_layout(slide_type, prefer_dark=prefer_dark)
        slide = self.prs.slides.add_slide(layout)

        # Handle speaker notes for all slide types
        if data.get("notes"):
            slide.notes_slide.notes_text_frame.text = data["notes"]

        return slide

    # =========================================================================
    # Direct Template Slides
    # =========================================================================

    def add_title_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create title slide using template layout.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("title", data, prefer_dark=prefer_dark)

        # Fill title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Presentation Title"))

        # Fill subtitle (idx 1) - may contain author/date
        subtitle_ph = self.get_placeholder(slide, idx=1)
        subtitle_parts = []
        if data.get("subtitle"):
            subtitle_parts.append(data["subtitle"])
        self.fill_text(subtitle_ph, "\n".join(subtitle_parts) if subtitle_parts else "")

        # Author/date in second subtitle if available (idx 2)
        author_ph = self.get_placeholder(slide, idx=2)
        if author_ph:
            author_parts = []
            if data.get("author"):
                author_parts.append(data["author"])
            if data.get("date"):
                author_parts.append(data["date"])
            self.fill_text(author_ph, " | ".join(author_parts) if author_parts else "")

    def add_section_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create section divider slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("section", data, prefer_dark=prefer_dark)

        # Fill title placeholder (idx 0) - template handles styling
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))

    def add_content_slide(self, data: Dict[str, Any]) -> None:
        """Create content slide with bullets."""
        slide = self._create_slide("content", data)

        # Title
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Slide Title"))

        # Subtitle (idx 2 for "7 Content A - Basic" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Bullets (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        self.fill_bullets(body_ph, data.get("bullets", []))

    def add_two_column_slide(self, data: Dict[str, Any]) -> None:
        """Create two-column slide."""
        slide = self._create_slide("two-column", data)

        # Title (type TITLE = 1)
        title_ph = self.get_placeholder(slide, ph_type=1)
        self.fill_text(title_ph, data.get("title", "Two Column"))

        # Subtitle row (idx 5 for "9 Content B - 2 Column" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Get SUBTITLE placeholders (type 4) sorted by position - these are column headers
        subtitle_phs = self.get_placeholders_by_type(slide, 4)
        # Column headers are at y > 1.5 inches (filter out top subtitle row)
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)  # Sort left to right

        # Get BODY placeholders (type 2) sorted by position - these are column content
        body_phs = self.get_placeholders_by_type(slide, 2)
        # Column bodies are at y > 2.5 inches
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)  # Sort left to right

        # Fill headers
        if data.get("left_header") and len(col_headers) > 0:
            self.fill_text(col_headers[0], data["left_header"])
        if data.get("right_header") and len(col_headers) > 1:
            self.fill_text(col_headers[1], data["right_header"])

        # Fill content
        if len(col_bodies) > 0:
            self.fill_bullets(col_bodies[0], data.get("left", []))
        if len(col_bodies) > 1:
            self.fill_bullets(col_bodies[1], data.get("right", []))

    def add_three_column_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column slide."""
        slide = self._create_slide("three-column", data)

        # Title (type TITLE = 1)
        title_ph = self.get_placeholder(slide, ph_type=1)
        self.fill_text(title_ph, data.get("title", "Three Column"))

        # Subtitle row (idx 7 for "11 Content C - 3 Column" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        columns = data.get("columns", [])

        # Get SUBTITLE placeholders (type 4) sorted by position - these are column headers
        # Filter out the top subtitle row (used for page subtitle) by checking y position
        subtitle_phs = self.get_placeholders_by_type(slide, 4)
        # Column headers are typically at y > 2 inches, subtitle row is at y < 2
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)  # Sort left to right

        # Get BODY placeholders (type 2) sorted by position - these are column content
        body_phs = self.get_placeholders_by_type(slide, 2)
        # Column bodies are typically at y > 2.5 inches
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)  # Sort left to right

        # Fill column headers
        for i, col in enumerate(columns[:3]):
            if i < len(col_headers) and col.get("header"):
                self.fill_text(col_headers[i], col["header"])

        # Fill column content
        for i, col in enumerate(columns[:3]):
            if i < len(col_bodies):
                self.fill_bullets(col_bodies[i], col.get("items", []))

    def add_big_number_slide(self, data: Dict[str, Any]) -> None:
        """Create big number/hero stat slide."""
        slide = self._create_slide("big-number", data)

        # Big number in title placeholder (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("number", "0"), bold=True, color=COLORS["accent"])

        # Description in body (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        description = data.get("text", "")
        if data.get("subtitle"):
            description += f"\n{data['subtitle']}"
        self.fill_text(body_ph, description)

    def add_callout_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create callout/bold statement slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("callout", data, prefer_dark=prefer_dark)

        # Main text in title placeholder (idx 0) - template handles styling
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("text", "Key message"))

        # Source attribution in subtitle placeholder (idx 1)
        if data.get("source"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['source']}")

    def add_quote_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create quote/testimonial slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("quote", data, prefer_dark=prefer_dark)

        # Quote text in title placeholder (idx 0) - template handles styling
        quote_text = f'"{data.get("quote", "Quote goes here.")}"'
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, quote_text)

        # Attribution in subtitle placeholder (idx 1)
        if data.get("attribution"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['attribution']}")

    def add_closing_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create closing/thank you slide.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("closing", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        text_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # The closing layout has pre-designed graphics (Databricks logo, etc.)
        title = data.get("title", "Thank You")

        # Title text - positioned at top
        self.add_textbox(slide, title, 0.75, 0.8, 11.5, 1.2,
                        font_size=48, bold=True, color=text_color,
                        alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # New Template-Based Slides (8 new types)
    # =========================================================================

    def add_two_column_icons_slide(self, data: Dict[str, Any]) -> None:
        """Create two-column slide with icon spots."""
        slide = self._create_slide("two-column-icons", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Two Column with Icons"))

        # Subtitle row (idx 5)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Headers (idx 3, 4)
        columns = data.get("columns", [])
        for i, col in enumerate(columns[:2]):
            header_ph = self.get_placeholder(slide, idx=3+i)
            if header_ph and col.get("header"):
                self.fill_text(header_ph, col["header"])

        # Content (idx 1, 2)
        for i, col in enumerate(columns[:2]):
            body_ph = self.get_placeholder(slide, idx=1+i)
            self.fill_bullets(body_ph, col.get("items", []))

        # Icons would go in the icon spots - template has picture placeholders
        # For now, users can add icons manually or we can extend this later

    def add_three_column_icons_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column slide with icon spots."""
        slide = self._create_slide("three-column-icons", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Three Column with Icons"))

        # Subtitle row (idx 7)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        columns = data.get("columns", [])

        # Headers (idx 3, 4, 6)
        header_indices = [3, 4, 6]
        for i, col in enumerate(columns[:3]):
            header_ph = self.get_placeholder(slide, idx=header_indices[i])
            if header_ph and col.get("header"):
                self.fill_text(header_ph, col["header"])

        # Content (idx 1, 2, 5)
        body_indices = [1, 2, 5]
        for i, col in enumerate(columns[:3]):
            body_ph = self.get_placeholder(slide, idx=body_indices[i])
            self.fill_bullets(body_ph, col.get("items", []))

    def add_cards_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column cards slide."""
        slide = self._create_slide("cards", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Cards"))

        # Subtitle row (idx 7)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        cards = data.get("cards", [])

        # Card headers (idx 4, 5, 6)
        for i, card in enumerate(cards[:3]):
            header_ph = self.get_placeholder(slide, idx=4+i)
            if header_ph and card.get("header"):
                self.fill_text(header_ph, card["header"])

        # Card content (idx 1, 2, 3)
        for i, card in enumerate(cards[:3]):
            body_ph = self.get_placeholder(slide, idx=1+i)
            if card.get("content"):
                self.fill_text(body_ph, card["content"])
            elif card.get("items"):
                self.fill_bullets(body_ph, card["items"])

    def add_card_right_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with content left, card right."""
        slide = self._create_slide("card-right", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Right"))

        # Subtitle (idx 3)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Left content (idx 1)
        left_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(left_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(left_ph, data["bullets"])

        # Right card area (idx 2) - for diagrams/images/tables
        right_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(right_ph, data["card_content"])

    def add_card_left_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with card left, content right."""
        slide = self._create_slide("card-left", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Left"))

        # Subtitle (idx 3)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Left card area (idx 2)
        left_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(left_ph, data["card_content"])

        # Right content (idx 1)
        right_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(right_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(right_ph, data["bullets"])

    def add_card_full_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with full-width card."""
        slide = self._create_slide("card-full", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Full Card"))

        # Subtitle (idx 2)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Card content (idx 1)
        card_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(card_ph, data["content"])

    def add_one_column_slide(self, data: Dict[str, Any]) -> None:
        """Create one-column text slide."""
        slide = self._create_slide("one-column", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", ""))

        # Subtitle row (idx 2 for "7 Content A - Basic" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Body (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(body_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(body_ph, data["bullets"])

    def add_section_description_slide(self, data: Dict[str, Any]) -> None:
        """Create section slide with description."""
        slide = self._create_slide("section-description", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))

        # Subtitle (idx 1)
        subtitle_ph = self.get_placeholder(slide, idx=1)
        if subtitle_ph:
            self.fill_text(subtitle_ph, data.get("subtitle", ""))

        # Body/description (idx 2)
        body_ph = self.get_placeholder(slide, idx=2)
        if data.get("description"):
            self.fill_text(body_ph, data["description"])
        elif data.get("bullets"):
            self.fill_bullets(body_ph, data["bullets"])

    # =========================================================================
    # Hybrid Slides (template background + custom shapes)
    # =========================================================================

    def add_agenda_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create agenda slide with custom hexagon numbers.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("agenda", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        item_bg_color = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Agenda"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Custom agenda items with hexagons
        items = data.get("items", [])
        start_y = 2.0

        for i, item in enumerate(items, 1):
            y_pos = start_y + (i - 1) * 0.9

            # Hexagon for number
            hex_shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON,
                Inches(0.75), Inches(y_pos),
                Inches(0.6), Inches(0.6)
            )
            hex_shape.fill.solid()
            hex_shape.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
            hex_shape.line.fill.background()

            # Number in hexagon (always light text on accent)
            tf = hex_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i)
            p.font.name = FONT_FAMILY
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Item text with background bar
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(8), Inches(0.6)
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = hex_to_rgb(item_bg_color)
            item_bg.line.fill.background()

            # Item text (always dark on light bar)
            self.add_textbox(slide, item, 1.7, y_pos + 0.1, 7.5, 0.5,
                           font_size=20, color=COLORS["text_dark"])

    def add_timeline_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create timeline/process slide with steps.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("timeline", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        secondary_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]

        # Title
        self.add_textbox(slide, data.get("title", "Timeline"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        steps = data.get("steps", [])
        num_steps = len(steps)
        if num_steps == 0:
            return

        step_width = 10.5 / num_steps
        start_x = 1.4

        # Connecting line (accent color works on both backgrounds)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x + 0.3), Inches(3.1),
            Inches(step_width * num_steps - 0.6), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
        line.line.fill.background()

        for i, step in enumerate(steps):
            x_pos = start_x + (i * step_width)

            # Circle for step number (accent fill, always light text)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + step_width/2 - 0.35), Inches(2.75),
                Inches(0.7), Inches(0.7)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
            circle.line.fill.background()

            # Step number (always light on accent)
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.name = FONT_FAMILY
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Step title
            self.add_textbox(slide, step.get("title", f"Step {i+1}"),
                           x_pos, 3.7, step_width, 0.6,
                           font_size=16, bold=True, color=body_color,
                           alignment=PP_ALIGN.CENTER)

            # Step description
            if step.get("description"):
                self.add_textbox(slide, step["description"],
                               x_pos, 4.4, step_width, 1.5,
                               font_size=12, color=secondary_color,
                               alignment=PP_ALIGN.CENTER)

    def add_icon_grid_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create icon grid slide for features/capabilities.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("icon-grid", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        secondary_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]
        circle_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Features"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        items = data.get("items", data.get("features", []))
        num_items = len(items)
        if num_items == 0:
            return

        # Determine grid layout
        if num_items <= 3:
            cols, rows = num_items, 1
        elif num_items <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 4, 2

        cell_width = 11 / cols
        cell_height = 2.2
        start_x = 1.2
        start_y = 1.8

        for i, item in enumerate(items[:8]):
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.5))

            # Circle with accent border
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + cell_width/2 - 0.5), Inches(y_pos),
                Inches(1), Inches(1)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(circle_fill)
            circle.line.color.rgb = hex_to_rgb(COLORS["accent"])
            circle.line.width = Pt(3)

            # Icon (emoji preferred, falls back to first letter of title)
            icon_raw = item.get("icon")
            if icon_raw:
                icon_text = icon_raw if len(icon_raw) <= 2 else icon_raw[0]
            else:
                icon_text = item.get("title", "?")[0].upper()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = icon_text
            p.font.name = FONT_FAMILY
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["accent"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Item title
            self.add_textbox(slide, item.get("title", ""),
                           x_pos, y_pos + 1.1, cell_width, 0.5,
                           font_size=14, bold=True, color=body_color,
                           alignment=PP_ALIGN.CENTER)

            # Item description
            if item.get("description"):
                self.add_textbox(slide, item["description"],
                               x_pos, y_pos + 1.55, cell_width, 0.8,
                               font_size=11, color=secondary_color,
                               alignment=PP_ALIGN.CENTER)

    def add_stat_row_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create slide with multiple stats in a row."""
        slide = self._create_slide("stat-row", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        label_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # Title
        self.add_textbox(slide, data.get("title", "Key Metrics"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        stats = data.get("stats", [])
        num_stats = len(stats)
        if num_stats == 0:
            return

        stat_width = 11.5 / num_stats
        start_x = 0.9

        for i, stat in enumerate(stats):
            x_pos = start_x + (i * stat_width)

            # Stat value (accent color works on both backgrounds)
            self.add_textbox(slide, stat.get("value", "0"),
                           x_pos, 2.5, stat_width - 0.3, 1.5,
                           font_size=56, bold=True, color=COLORS["accent"],
                           alignment=PP_ALIGN.CENTER)

            # Stat label
            self.add_textbox(slide, stat.get("label", ""),
                           x_pos, 4.2, stat_width - 0.3, 1.0,
                           font_size=16, bold=True, color=label_color,
                           alignment=PP_ALIGN.CENTER)

            # Divider (except after last)
            if i < num_stats - 1:
                divider = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos + stat_width - 0.15), Inches(2.7),
                    Inches(0.02), Inches(2.5)
                )
                divider.fill.solid()
                divider.fill.fore_color.rgb = hex_to_rgb(COLORS["divider"])
                divider.line.fill.background()

    def add_pros_cons_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create pros/cons comparison slide."""
        slide = self._create_slide("pros-cons", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # Title
        self.add_textbox(slide, data.get("title", "Pros & Cons"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Pros header (green works on both backgrounds)
        self.add_textbox(slide, data.get("pros_header", "Pros"),
                        0.75, 1.6, 5.5, 0.5,
                        font_size=20, bold=True, color=COLORS["green"])

        # Pros items
        pros = data.get("pros", [])
        for i, pro in enumerate(pros):
            self.add_textbox(slide, f"✓  {pro}",
                           0.75, 2.2 + (i * 0.6), 5.5, 0.5,
                           font_size=16, color=body_color)

        # Cons header (red works on both backgrounds)
        self.add_textbox(slide, data.get("cons_header", "Cons"),
                        7.0, 1.6, 5.5, 0.5,
                        font_size=20, bold=True, color=COLORS["red"])

        # Cons items
        cons = data.get("cons", [])
        for i, con in enumerate(cons):
            self.add_textbox(slide, f"✗  {con}",
                           7.0, 2.2 + (i * 0.6), 5.5, 0.5,
                           font_size=16, color=body_color)

    def add_comparison_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create VS comparison slide."""
        slide = self._create_slide("comparison", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        label_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        diamond_bg = COLORS["dark_bg"] if prefer_dark else COLORS["accent"]
        diamond_text = COLORS["text_light"]  # Always white on diamond

        # Title
        self.add_textbox(slide, data.get("title", "Comparison"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # VS diamond in center
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(6.166), Inches(3.25),
            Inches(1), Inches(1)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = hex_to_rgb(diamond_bg)
        diamond.line.fill.background()

        # VS text
        tf = diamond.text_frame
        p = tf.paragraphs[0]
        p.text = "vs."
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(diamond_text)
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Left label
        self.add_textbox(slide, data.get("left_label", "Option A"),
                        1.5, 5.0, 4.0, 0.6,
                        font_size=20, bold=True, color=label_color,
                        alignment=PP_ALIGN.CENTER)

        # Right label
        self.add_textbox(slide, data.get("right_label", "Option B"),
                        7.833, 5.0, 4.0, 0.6,
                        font_size=20, bold=True, color=label_color,
                        alignment=PP_ALIGN.CENTER)

    def add_checklist_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create checklist slide."""
        slide = self._create_slide("checklist", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        unchecked_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Checklist"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        items = data.get("items", [])
        start_y = 1.8

        for i, item in enumerate(items):
            y_pos = start_y + (i * 0.7)

            # Checkbox
            checkbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.9), Inches(y_pos),
                Inches(0.35), Inches(0.35)
            )
            checkbox.fill.solid()

            is_checked = item.get("checked", False) if isinstance(item, dict) else False
            item_text = item.get("text", item) if isinstance(item, dict) else item

            if is_checked:
                checkbox.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
                # Checkmark (always white on accent background)
                tf = checkbox.text_frame
                p = tf.paragraphs[0]
                p.text = "✓"
                p.font.name = FONT_FAMILY
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
                p.alignment = PP_ALIGN.CENTER
                tf.anchor = MSO_ANCHOR.MIDDLE
            else:
                checkbox.fill.fore_color.rgb = hex_to_rgb(unchecked_fill)

            checkbox.line.color.rgb = hex_to_rgb(COLORS["accent"])
            checkbox.line.width = Pt(2)

            # Item text
            self.add_textbox(slide, item_text, 1.5, y_pos, 10.0, 0.4,
                           font_size=16, color=body_color)

    def add_logos_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create logo display slide."""
        slide = self._create_slide("logos", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        subtitle_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]

        # Title
        self.add_textbox(slide, data.get("title", "Our Partners"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Subtitle
        if data.get("subtitle"):
            self.add_textbox(slide, data["subtitle"],
                           0.75, 1.3, 11.0, 0.5,
                           font_size=16, color=subtitle_color,
                           alignment=PP_ALIGN.CENTER)

        logos = data.get("logos", [])
        num_logos = len(logos)

        # Grid layout
        if num_logos <= 4:
            cols, rows = num_logos, 1
        elif num_logos <= 8:
            cols, rows = 4, 2
        else:
            cols, rows = 5, 2

        cell_width = 10 / cols
        cell_height = 1.6
        start_x = 1.7
        start_y = 2.5

        # Logo box colors based on background
        box_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]
        box_text = COLORS["text_secondary"]  # Gray works on white boxes

        for i, logo in enumerate(logos[:10]):
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.3))

            # Logo placeholder box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(cell_width - 0.4), Inches(cell_height - 0.3)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = hex_to_rgb(box_fill)
            box.line.color.rgb = hex_to_rgb(COLORS["divider"])
            box.line.width = Pt(1)

            # Company name as placeholder
            logo_name = logo if isinstance(logo, str) else logo.get("name", "Company")
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = logo_name
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(box_text)
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

    # =========================================================================
    # Generation Methods
    # =========================================================================

    def generate(self, content: Dict[str, Any]) -> None:
        """Generate all slides from content dict."""
        slides = content.get("slides", [])

        slide_methods = {
            # Direct template slides
            "title": self.add_title_slide,
            "section": self.add_section_slide,
            "content": self.add_content_slide,
            "two-column": self.add_two_column_slide,
            "three-column": self.add_three_column_slide,
            "big-number": self.add_big_number_slide,
            "callout": self.add_callout_slide,
            "quote": self.add_quote_slide,
            "closing": self.add_closing_slide,
            # New template slides
            "two-column-icons": self.add_two_column_icons_slide,
            "three-column-icons": self.add_three_column_icons_slide,
            "cards": self.add_cards_slide,
            "card-right": self.add_card_right_slide,
            "card-left": self.add_card_left_slide,
            "card-full": self.add_card_full_slide,
            "one-column": self.add_one_column_slide,
            "section-description": self.add_section_description_slide,
            # Hybrid slides
            "agenda": self.add_agenda_slide,
            "timeline": self.add_timeline_slide,
            "icon-grid": self.add_icon_grid_slide,
            "stat-row": self.add_stat_row_slide,
            "pros-cons": self.add_pros_cons_slide,
            "comparison": self.add_comparison_slide,
            "checklist": self.add_checklist_slide,
            "logos": self.add_logos_slide,
        }

        for slide_data in slides:
            slide_type = slide_data.get("type", "content")
            if slide_type not in VALID_SLIDE_TYPES:
                print(f"Warning: Unknown slide type '{slide_type}', using 'content'")
                slide_type = "content"
            method = slide_methods.get(slide_type, self.add_content_slide)
            method(slide_data)

    def save(self, output_path: str) -> str:
        """Save presentation to file."""
        self.prs.save(output_path)
        return output_path


# =============================================================================
# CLI
# =============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate Databricks-branded PowerPoint presentations"
    )
    parser.add_argument("--input", "-i", required=True, help="Path to JSON content file")
    parser.add_argument("--output", "-o", required=True, help="Output path for .pptx file")

    args = parser.parse_args()

    # Load content
    with open(args.input, encoding='utf-8') as f:
        content = json.load(f)

    # Generate presentation
    generator = DatabricksSlideGenerator()
    generator.generate(content)
    output_path = generator.save(args.output)

    print(f"✓ Generated: {output_path}")
    print(f"  Slides: {generator.slide_count}")


if __name__ == "__main__":
    main()
